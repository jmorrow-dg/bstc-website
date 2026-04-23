#!/usr/bin/env python3
"""Generate How I AI Edition 3 PPTX.

Style mirrors Edition 1/2: charcoal background, red accent, white display text,
muted grey body. No em dashes or en dashes. Australian English.

Slide order (page numbers slot into a 15-slide full deck):
  1. Jargon Buster             (05/15)
  2. Signal 01, Vercel Breach  (06/15)
  3. Signal 02, Prompt Injection (07/15)
  4. Signal 03, Claude Code Pricing (08/15)
  5. Signal 04, Thunderbolt    (09/15)
  6. Signal 05, Stanford Index (10/15)
  7. Action of the Week        (11/15)
  8. By Tier                   (12/15)
  9. Community Wins            (15/15)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

CHARCOAL = RGBColor(0x11, 0x11, 0x11)
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
RED = RGBColor(0xC8, 0x1E, 0x1E)
GOLD = RGBColor(0xE0, 0xB0, 0x40)
GREY = RGBColor(0x88, 0x88, 0x88)
GREY_LIGHT = RGBColor(0xAA, 0xAA, 0xAA)
GREY_DARK = RGBColor(0x33, 0x33, 0x33)
GREY_MID = RGBColor(0x44, 0x44, 0x44)
GREY_BORDER = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = CHARCOAL


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                font_name="Arial"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return tb


def add_multiline(slide, left, top, width, height, lines, default_size=18,
                  default_color=WHITE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.get("text", "")
        p.font.size = Pt(line.get("size", default_size))
        p.font.color.rgb = line.get("color", default_color)
        p.font.bold = line.get("bold", False)
        p.font.italic = line.get("italic", False)
        p.font.name = line.get("font", "Arial")
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        if line.get("spacing_before"):
            p.space_before = Pt(line["spacing_before"])
    return tb


def add_tag(slide, left, top, text):
    add_textbox(slide, left, top, Inches(8), Inches(0.4), text,
                font_size=13, color=RED, bold=True)


def add_accent_line(slide, left, top, width=Inches(0.8), color=RED, thickness=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, border_color=GREY_BORDER,
             fill_color=GREY_DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape


def add_red_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color=GREY_MID, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    return shape


def add_footer(slide, num, total=15):
    add_textbox(slide, Inches(0.5), Inches(7.1), Inches(1.5), Inches(0.3),
                "BSTC", font_size=11, color=GREY_DARK, bold=True)
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
                f"{num:02d} / {total}", font_size=12, color=GREY,
                alignment=PP_ALIGN.RIGHT)


def add_takeaway(slide, left, top, width, text, height=Inches(1.2)):
    """'What this means for you' card with red left bar."""
    add_card(slide, left, top, width, height)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    add_multiline(slide, left + Inches(0.2), top + Inches(0.12),
                  width - Inches(0.35), height - Inches(0.2), [
        {"text": "What this means for you:", "size": 13, "color": WHITE, "bold": True},
        {"text": text, "size": 12, "color": GREY_LIGHT, "spacing_before": 6},
    ])


def add_date_tag(slide, left, top, text):
    add_textbox(slide, left, top, Inches(8), Inches(0.3), text,
                font_size=11, color=GREY, bold=True)


def add_tier_card(slide, x, y, width, height, label, heading, bullets, goal, highlight=False):
    """Card with red top border, red label, white heading, grey bullets, red goal."""
    add_card(slide, x, y, width, height,
             border_color=GREY_BORDER, fill_color=GREY_DARK)
    # Red top border strip
    add_accent_line(slide, x, y, width=width, thickness=Pt(4))
    add_textbox(slide, x + Inches(0.3), y + Inches(0.25), width - Inches(0.6), Inches(0.3),
                label, font_size=12, color=RED, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.65), width - Inches(0.6), Inches(0.6),
                heading, font_size=20, color=WHITE, bold=True)
    items_text = "\n\n".join(f"\u2022  {item}" for item in bullets)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.5), width - Inches(0.6), height - Inches(2.3),
                items_text, font_size=12, color=GREY_LIGHT)
    add_textbox(slide, x + Inches(0.3), y + height - Inches(0.7), width - Inches(0.6), Inches(0.4),
                goal, font_size=12, color=RED, bold=True, italic=True)


# ============================================================
# SLIDE: Jargon Buster  (05/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "BEFORE WE START")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.9),
            "Jargon Buster", font_size=48, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(2.15))
add_textbox(slide, Inches(0.8), Inches(2.45), Inches(11), Inches(0.6),
            "Five terms you will hear tonight. Lock these in first and the signals land harder.",
            font_size=16, color=GREY_LIGHT)

jargon = [
    ("OAuth",
     "The permission slip you sign when an app asks 'continue with Google.' It can stay valid for years."),
    ("MCP  (Model Context Protocol)",
     "The USB-C of AI tools. Lets agents plug into your systems in a standard way."),
    ("Non-human identity",
     "A user account that is not a person. Your scripts, bots, and agents all have them."),
    ("Prompt injection",
     "Hiding instructions inside data an AI reads, so it executes commands you did not send."),
    ("Agent",
     "AI that takes actions on your behalf, not just answers questions."),
]

row_top = Inches(3.3)
row_h = Inches(0.62)
term_left = Inches(0.8)
term_w = Inches(3.8)
def_left = Inches(4.8)
def_w = Inches(7.9)

for i, (term, definition) in enumerate(jargon):
    y = row_top + row_h * i
    add_textbox(slide, term_left, y + Inches(0.08), term_w, row_h,
                term, font_size=16, color=RED, bold=True)
    add_textbox(slide, def_left, y + Inches(0.08), def_w, row_h,
                definition, font_size=14, color=WHITE)
    # Gold divider at bottom of each row
    add_accent_line(slide, term_left, y + row_h - Inches(0.02),
                    width=Inches(11.9), color=GOLD, thickness=Pt(1))

add_textbox(slide, Inches(0.8), Inches(6.85), Inches(11.9), Inches(0.3),
            "If any of these feel fuzzy, find me at dinner. No dumb questions in this room.",
            font_size=13, color=GREY, italic=True)

add_footer(slide, 5)


# ============================================================
# SIGNAL 01: The Vercel Breach  (06/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "SIGNAL 01")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.9),
            "The Vercel Breach", font_size=48, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(2.2))
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.0),
            "Attackers compromised Vercel through an AI agent platform's OAuth scope. "
            "No model exploit. No prompt injection. OAuth was the skeleton key. "
            "Customer environment variables were exposed.",
            font_size=16, color=GREY_LIGHT)
add_date_tag(slide, Inches(0.8), Inches(3.45), "19 APRIL 2026")

chain = [
    ("Context.ai", "AI agent platform", GREY_DARK, WHITE, GREY_BORDER),
    ("OAuth scope", "Google Workspace", GREY_DARK, GOLD, GOLD),
    ("Employee account", "Vercel", GREY_DARK, WHITE, GREY_BORDER),
    ("Internal systems", "Customer env vars", GREY_DARK, WHITE, GREY_BORDER),
]
box_w = Inches(2.7)
box_h = Inches(1.1)
start_x = Inches(0.6)
gap = Inches(0.25)
chain_top = Inches(3.95)
for i, (title, sub, fill, title_color, border) in enumerate(chain):
    x = start_x + (box_w + gap) * i
    add_card(slide, x, chain_top, box_w, box_h, border_color=border, fill_color=fill)
    add_multiline(slide, x, chain_top + Inches(0.2), box_w, Inches(0.9), [
        {"text": title, "size": 14, "color": title_color, "bold": True, "align": PP_ALIGN.CENTER},
        {"text": sub, "size": 11, "color": GREY_LIGHT, "align": PP_ALIGN.CENTER, "spacing_before": 6},
    ])
    if i < 3:
        arrow_x = x + box_w + Inches(0.02)
        arrow_color = GOLD if i == 1 else RED
        add_textbox(slide, arrow_x, chain_top + Inches(0.35), gap, Inches(0.5),
                    "\u2192", font_size=20, color=arrow_color,
                    alignment=PP_ALIGN.CENTER, bold=True)

add_textbox(slide, Inches(0.8), Inches(5.25), Inches(11.7), Inches(0.4),
            "No model exploited. No prompt injection. OAuth scope was the skeleton key.",
            font_size=13, color=GOLD, italic=True, alignment=PP_ALIGN.CENTER)

add_takeaway(slide, Inches(0.8), Inches(5.85), Inches(11.7),
             "Every AI agent you install asks for OAuth scopes into Gmail, Drive, Slack, and "
             "your workspace. Each token is a potential skeleton key to your infrastructure. "
             "Review what you have connected, revoke what you do not use, and rotate long-lived "
             "secrets. The agent hype cycle just met its first real-world bill.",
             height=Inches(1.15))

add_footer(slide, 6)


# ============================================================
# SIGNAL 02: Prompt Injection & Agent Identity  (07/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "SIGNAL 02")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(2.0),
            "Prompt Injection and\nAgent Identity Named\nTop Enterprise Threat",
            font_size=36, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(3.6))
add_textbox(slide, Inches(0.8), Inches(3.9), Inches(7.5), Inches(1.2),
            "Security researchers flagged non-human identity and prompt injection as the "
            "leading enterprise AI threats five days before Vercel got hit through exactly "
            "that class of attack.",
            font_size=15, color=GREY_LIGHT)
add_date_tag(slide, Inches(0.8), Inches(5.3), "14 APRIL 2026  \u00b7  ISACA")

add_takeaway(slide, Inches(0.8), Inches(5.85), Inches(7.5),
             "Signal 01 is Signal 02 manifesting in production. If you are not governing "
             "agent identity, you are running the risk Vercel just learned the hard way. "
             "Start with OAuth scope hygiene, short-lived tokens, and agent permission scoping.",
             height=Inches(1.15))

stats_2 = [
    ("100 : 1", "Machine-to-machine vs human logins in some enterprises"),
    ("Attack surface", "GitHub issues, docs, email, PR descriptions"),
    ("Blind spot", "Existing endpoint and IAM tooling does not cover service principals"),
]
card_w = Inches(4.3)
card_h = Inches(1.5)
right_x = Inches(8.6)
for i, (headline, caption) in enumerate(stats_2):
    y = Inches(1.3) + (card_h + Inches(0.2)) * i
    add_card(slide, right_x, y, card_w, card_h)
    add_multiline(slide, right_x + Inches(0.2), y + Inches(0.2),
                  card_w - Inches(0.4), card_h - Inches(0.3), [
        {"text": headline, "size": 26, "color": RED, "bold": True, "align": PP_ALIGN.LEFT},
        {"text": caption, "size": 12, "color": GREY_LIGHT, "spacing_before": 6},
    ])

add_footer(slide, 7)


# ============================================================
# SIGNAL 03: Claude Code Pricing  (08/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "SIGNAL 03")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(1.8),
            "The End of\nSubsidised AI", font_size=44, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(3.3))
add_textbox(slide, Inches(0.8), Inches(3.6), Inches(7.5), Inches(1.5),
            "Anthropic quietly pulled Claude Code from the $20 Pro plan, routing users to the "
            "$100 Max tier. Official framing: a 2% test on new signups. Reality: support docs "
            "updated globally, with no announcement and no grandfathering for paying users.",
            font_size=15, color=GREY_LIGHT)
add_date_tag(slide, Inches(0.8), Inches(5.2), "21 APRIL 2026  \u00b7  CLAUDE CODE REMOVED FROM $20 PRO TIER")

add_takeaway(slide, Inches(0.8), Inches(5.6), Inches(7.5),
             "The VC funded token subsidy era is ending in public. If your product economics "
             "assume a flat monthly ceiling from Anthropic, OpenAI, or a reseller, re-stress-test "
             "margins for token pass-through within 12 months. For SEA teams, $100 USD parity "
             "bites harder. Expect faster shift to GLM, Kimi, Qwen, Cursor, and Codex. Framing: "
             "the companies that hid the subsidy best will eat the most backlash.",
             height=Inches(1.45))

stats_3 = [
    ("5x", "Minimum subscription jump"),
    ("10x", "Token cost vs plan fee, power user"),
    ("2%", "Anthropic 2% test cohort claim"),
]
for i, (big, caption) in enumerate(stats_3):
    y = Inches(1.3) + (Inches(1.5) + Inches(0.2)) * i
    add_card(slide, Inches(8.6), y, Inches(4.3), Inches(1.5))
    add_multiline(slide, Inches(8.6), y + Inches(0.18),
                  Inches(4.3), Inches(1.3), [
        {"text": big, "size": 44, "color": RED, "bold": True, "align": PP_ALIGN.CENTER},
        {"text": caption, "size": 12, "color": GREY_LIGHT, "align": PP_ALIGN.CENTER, "spacing_before": 4},
    ])

add_footer(slide, 8)


# ============================================================
# SIGNAL 04: Mozilla Thunderbolt  (09/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "SIGNAL 04")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(1.4),
            "Mozilla Thunderbolt\nLaunches", font_size=44, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(3.3))
add_textbox(slide, Inches(0.8), Inches(3.6), Inches(7.5), Inches(1.5),
            "Mozilla shipped an open-source, self-hostable enterprise AI client as a direct "
            "alternative to Microsoft Copilot, ChatGPT Enterprise, and Claude Enterprise. "
            "Any model, MCP native, end-to-end encrypted.",
            font_size=15, color=GREY_LIGHT)
add_date_tag(slide, Inches(0.8), Inches(5.3), "16 APRIL 2026")

add_takeaway(slide, Inches(0.8), Inches(5.85), Inches(7.5),
             "Ties directly to Signal 01. If your data never left your own infrastructure, "
             "a Vercel-style third-party breach would not touch you. The objection of 'we "
             "cannot use AI because of data concerns' just died. Strategic question: which of "
             "your workflows warrant self-hosted AI?",
             height=Inches(1.15))

features = [
    "Any model: commercial, open-source, locally hosted",
    "MCP and Agent Client Protocol native",
    "End-to-end encryption, device-level access control",
    "Available on GitHub, managed version for small teams",
]
for i, feat in enumerate(features):
    y = Inches(1.3) + (Inches(1.1) + Inches(0.15)) * i
    add_card(slide, Inches(8.6), y, Inches(4.3), Inches(1.1))
    add_textbox(slide, Inches(8.85), y + Inches(0.22), Inches(0.35), Inches(0.5),
                "\u25B6", font_size=14, color=RED, bold=True)
    add_textbox(slide, Inches(9.25), y + Inches(0.18), Inches(3.5), Inches(0.9),
                feat, font_size=13, color=WHITE, bold=True)

add_footer(slide, 9)


# ============================================================
# SIGNAL 05: Stanford AI Index 2026  (10/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "SIGNAL 05")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(1.4),
            "Stanford AI Index 2026", font_size=44, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(2.4))
add_textbox(slide, Inches(0.8), Inches(2.7), Inches(7.5), Inches(1.5),
            "AI agent task success rate jumped from 20% to 77% in twelve months. GenAI reached "
            "53% population adoption in three years, the fastest in tech history.",
            font_size=15, color=GREY_LIGHT)
add_date_tag(slide, Inches(0.8), Inches(4.4), "15 APRIL 2026")

add_takeaway(slide, Inches(0.8), Inches(5.85), Inches(7.5),
             "If you tried agents six months ago and they felt unreliable, try again. The "
             "reliability gap just closed. This is the reason every SaaS is rushing to ship "
             "agent toolkits and why Signal 01 is now possible. The capability is ready. The "
             "security posture around it is not.",
             height=Inches(1.15))

left_sq_x = Inches(8.6)
arrow_x = Inches(10.25)
right_sq_x = Inches(10.85)
sq_top = Inches(1.6)
left_sq_size = Inches(1.5)
right_sq_size = Inches(2.0)

add_card(slide, left_sq_x, sq_top + Inches(0.25), left_sq_size, left_sq_size)
add_multiline(slide, left_sq_x, sq_top + Inches(0.45), left_sq_size, left_sq_size, [
    {"text": "20%", "size": 32, "color": GREY_LIGHT, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "2025", "size": 11, "color": GREY, "align": PP_ALIGN.CENTER, "spacing_before": 6},
])

add_textbox(slide, arrow_x, sq_top + Inches(0.85), Inches(0.6), Inches(0.5),
            "\u2192", font_size=28, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

add_red_card(slide, right_sq_x, sq_top, right_sq_size, right_sq_size)
add_multiline(slide, right_sq_x, sq_top + Inches(0.4), right_sq_size, right_sq_size, [
    {"text": "77%", "size": 50, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "2026", "size": 12, "color": WHITE, "align": PP_ALIGN.CENTER, "spacing_before": 6},
])

add_card(slide, Inches(8.6), Inches(4.0), Inches(4.3), Inches(1.2))
add_multiline(slide, Inches(8.6), Inches(4.25), Inches(4.3), Inches(0.9), [
    {"text": "53%", "size": 30, "color": RED, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "Population adoption in 3 years", "size": 12, "color": GREY_LIGHT,
     "align": PP_ALIGN.CENTER, "spacing_before": 4},
])

add_footer(slide, 10)


# ============================================================
# SLIDE: Action of the Week  (11/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "ACTION OF THE WEEK")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1.0),
            "Audit Your OAuth Grants Tonight", font_size=40, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(2.2))
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(9.5), Inches(1.0),
            "The single highest-leverage thing you can do before next Wednesday. "
            "Ten minutes, measurable risk reduction, proof you heard Signal 01.",
            font_size=15, color=GREY_LIGHT)

# 10 MINUTES badge, top right
badge_x = Inches(10.5)
badge_y = Inches(1.2)
badge_w = Inches(2.4)
badge_h = Inches(1.8)
add_red_card(slide, badge_x, badge_y, badge_w, badge_h)
add_multiline(slide, badge_x, badge_y + Inches(0.35), badge_w, badge_h - Inches(0.4), [
    {"text": "10", "size": 56, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER},
    {"text": "MINUTES", "size": 16, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER, "spacing_before": 2},
    {"text": "Average time to complete", "size": 10, "color": WHITE,
     "align": PP_ALIGN.CENTER, "spacing_before": 6},
])

# Four-step horizontal flow
steps = [
    ("01", "Open", "myaccount.google.com/permissions"),
    ("02", "Scan", "Every app listed. Ask: do I still use this weekly?"),
    ("03", "Revoke", "Anything unused in the last 30 days"),
    ("04", "Repeat", "For Slack, GitHub, Notion, Linear"),
]
flow_top = Inches(3.7)
flow_h = Inches(1.7)
step_w = Inches(2.95)
step_gap = Inches(0.1)
flow_start_x = Inches(0.6)
for i, (num, title, body) in enumerate(steps):
    x = flow_start_x + (step_w + step_gap) * i
    add_card(slide, x, flow_top, step_w, flow_h)
    # Red numbered circle
    circle_size = Inches(0.6)
    add_circle(slide, x + Inches(0.3), flow_top + Inches(0.2), circle_size, fill_color=RED)
    add_textbox(slide, x + Inches(0.3), flow_top + Inches(0.28), circle_size, Inches(0.5),
                num, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Step title
    add_textbox(slide, x + Inches(1.0), flow_top + Inches(0.28), step_w - Inches(1.1), Inches(0.45),
                title, font_size=18, color=WHITE, bold=True)
    # Body
    add_textbox(slide, x + Inches(0.3), flow_top + Inches(0.95), step_w - Inches(0.5), Inches(0.7),
                body, font_size=12, color=GREY_LIGHT)
    # Arrow between cards
    if i < 3:
        ax = x + step_w + Inches(0.005)
        add_textbox(slide, ax - Inches(0.05), flow_top + Inches(0.6), step_gap + Inches(0.1), Inches(0.5),
                    "\u203A", font_size=22, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

add_takeaway(slide, Inches(0.8), Inches(5.75), Inches(11.7),
             "Vercel got breached through an OAuth token to a product most employees had "
             "forgotten they installed. You almost certainly have the same kind of forgotten "
             "grant sitting in your Google account right now. Revoke first, ask questions later. "
             "If you need the permission back, the reconnect takes thirty seconds.",
             height=Inches(1.25))

add_footer(slide, 11)


# ============================================================
# SLIDE: By Tier  (12/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "BY TIER")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.9),
            "What To Do With This Week's News", font_size=40, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(2.2))
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(12), Inches(0.6),
            "One action calibrated to where you are right now. Pick your column.",
            font_size=16, color=GREY_LIGHT)

tiers = [
    {
        "label": "TIER 1, EXPLORER",
        "heading": "Learn the vocabulary",
        "bullets": [
            "Open Claude free tier this week",
            "Ask it: 'Explain OAuth permissions to me like I am signing up for my first app'",
            "Then ask it to help you audit one grant on your Google account",
        ],
        "goal": "Goal: turn abstract risk into concrete understanding",
    },
    {
        "label": "TIER 2, OPERATOR",
        "heading": "Audit your agent stack",
        "bullets": [
            "List every AI tool you have connected: ChatGPT plugins, Claude connectors, Zapier AI, Notion AI, browser extensions",
            "For each one, read the permission scope it was granted",
            "Revoke anything you do not use weekly",
        ],
        "goal": "Goal: know what your AI can reach before attackers do",
    },
    {
        "label": "TIER 3, BUILDER",
        "heading": "Ship secure agents",
        "bullets": [
            "If your product issues OAuth tokens, enforce short-lived expiry (1 hour, not 1 year)",
            "Scope permissions to the minimum your agent needs, not everything the API allows",
            "Log every agent action with the identity that triggered it",
        ],
        "goal": "Goal: do not be the next Vercel case study",
    },
]
tier_w = Inches(4.0)
tier_h = Inches(4.1)
tier_top = Inches(3.3)
tier_gap = Inches(0.15)
tier_start = Inches(0.5)
for i, t in enumerate(tiers):
    x = tier_start + (tier_w + tier_gap) * i
    add_tier_card(slide, x, tier_top, tier_w, tier_h,
                  t["label"], t["heading"], t["bullets"], t["goal"])

add_footer(slide, 12)


# ============================================================
# SLIDE: Community Wins  (15/15)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, Inches(0.8), Inches(0.6), "COMMUNITY WINS")
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(12), Inches(0.9),
            "Wins From The Room", font_size=44, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(2.2))
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(12), Inches(0.8),
            "Three members shipped something since last edition. This is what the community "
            "looks like when it builds in public.",
            font_size=15, color=GREY_LIGHT)

wins = [
    {
        "name": "[Member 1 name]",
        "role": "[Founder, Company  /  Operator, Industry]",
        "shipped": "Built an agent that drafts sales follow-ups from Granola transcripts. "
                   "40 minutes back per deal.",
        "initial": "M1",
    },
    {
        "name": "[Member 2 name]",
        "role": "[Founder or operator title]",
        "shipped": "Replaced their customer support tier-1 with Claude plus three MCP "
                   "connectors. Response time dropped from 4 hours to 8 minutes.",
        "initial": "M2",
    },
    {
        "name": "[Member 3 name]",
        "role": "[Founder or operator title]",
        "shipped": "Built an internal AI brief generator that runs every morning at 6am. "
                   "Their team starts the day aligned instead of in standups.",
        "initial": "M3",
    },
]
card_w = Inches(4.0)
card_h = Inches(3.4)
card_top = Inches(3.5)
card_gap = Inches(0.15)
card_start = Inches(0.5)
for i, w in enumerate(wins):
    x = card_start + (card_w + card_gap) * i
    add_card(slide, x, card_top, card_w, card_h)
    # Avatar circle
    avatar_size = Inches(1.1)
    avatar_x = x + (card_w - avatar_size) / 2
    add_circle(slide, avatar_x, card_top + Inches(0.3), avatar_size,
               fill_color=GREY_MID, border_color=RED)
    add_textbox(slide, avatar_x, card_top + Inches(0.58), avatar_size, Inches(0.6),
                w["initial"], font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Name
    add_textbox(slide, x + Inches(0.3), card_top + Inches(1.55), card_w - Inches(0.6), Inches(0.4),
                w["name"], font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Role
    add_textbox(slide, x + Inches(0.3), card_top + Inches(2.0), card_w - Inches(0.6), Inches(0.3),
                w["role"], font_size=12, color=GREY, italic=True, alignment=PP_ALIGN.CENTER)
    # Shipped
    add_textbox(slide, x + Inches(0.3), card_top + Inches(2.45), card_w - Inches(0.6), Inches(0.9),
                w["shipped"], font_size=12, color=GREY_LIGHT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(7.1), Inches(11.9), Inches(0.35),
            "Want to be on this slide next edition? Tell me what you shipped. "
            "No pitch needed, just a screenshot and one line.",
            font_size=13, color=GOLD, italic=True, alignment=PP_ALIGN.CENTER)

add_footer(slide, 15)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/joshmorrowdavidgoliath/bstc-website/public/how-i-ai-edition-3.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
